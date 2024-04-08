import streamlit as st
import nltk
from nltk.corpus import stopwords
from nltk.stem import PorterStemmer
from PyPDF2 import PdfReader
import docx
from pptx import Presentation
import re
from sklearn.feature_extraction.text import TfidfVectorizer

nltk.download('stopwords')
stop_words = stopwords.words('english')

try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')

stemmer = PorterStemmer()

def extract_keywords(text):
    vectorizer = TfidfVectorizer(stop_words=stop_words, ngram_range=(1, 2))
    tfidf_matrix = vectorizer.fit_transform([text])
    features = vectorizer.get_feature_names_out()
    scores = tfidf_matrix.toarray()[0]
    keywords = [(stemmer.stem(features[idx]), scores[idx]) for idx in range(len(features))]
    return keywords

def clean_and_filter_keywords(keywords):
    if not keywords:
        return []

    unique_keywords = set()
    filtered_keywords = []

    for keyword, score in keywords:
        keyword_cleaned = re.sub(r'[^a-zA-Z0-9\s]', '', keyword).lower()

        if len(keyword_cleaned) > 4 and keyword_cleaned not in stop_words:
            if len(keyword_cleaned.split()) <= 2:
                if keyword_cleaned not in unique_keywords:
                    unique_keywords.add(keyword_cleaned)
                    filtered_keywords.append((keyword_cleaned, score))

    return filtered_keywords


def extract_text_from_pdf(uploaded_file):
    pdf_reader = PdfReader(uploaded_file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_docx(uploaded_file):
    doc = docx.Document(uploaded_file)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + " "
    return text

def extract_text_from_pptx(uploaded_file):
    prs = Presentation(uploaded_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    return text


def main():
    st.set_page_config(page_title="Keyword Extraction App Using TF-IDF", layout="wide")
    st.title("Keyword Extraction App Using TF-IDF")

    col1, col2 = st.columns([2, 3])

    with col1:
        st.subheader("Upload File")
        uploaded_file = st.file_uploader("Choose a file (PDF, DOCX, PPTX, or TXT)", type=["pdf", "docx", "pptx", "txt"])

        if uploaded_file is not None:
            file_extension = uploaded_file.name.split(".")[-1]

            if file_extension == "pdf":
                with st.spinner("Extracting text from PDF..."):
                    text = extract_text_from_pdf(uploaded_file)
            elif file_extension == "docx":
                with st.spinner("Extracting text from DOCX..."):
                    text = extract_text_from_docx(uploaded_file)
            elif file_extension == "pptx":
                with st.spinner("Extracting text from PPTX..."):
                    text = extract_text_from_pptx(uploaded_file)
            elif file_extension == "txt":
                with st.spinner("Reading text file..."):
                    text = uploaded_file.read().decode('utf-8')
            else:
                st.error("Unsupported file type. Please upload a PDF, DOCX, PPTX, or TXT file.")
                return

            with st.spinner("Extracting keywords..."):
                keywords = extract_keywords(text)
                filtered_keywords = clean_and_filter_keywords(keywords)

            st.subheader("File Contents")
            if len(text) > 2500:
                show_more = st.checkbox("Show More", value=False)
                if show_more:
                    st.write(text)
                else:
                    st.write(text[:2000] + "...")
            else:
                st.write(text)

    with col2:
        if uploaded_file is not None:
            st.subheader("Extracted Keywords")
            top100 = st.checkbox("Top 100 Keywords", value=True)
            keyword_filter = st.text_input("Filter keywords (leave blank to show all)").lower()
            filtered_keywords = [kw for kw in filtered_keywords if keyword_filter in kw[0]]
            filtered_keywords.sort(key=lambda x: x[1], reverse=True)
            
            if top100:
                filtered_keywords = filtered_keywords[:100]
            
            if st.button("Download Keywords as Text File"):
                keywords_content = "\n".join([f"{keyword} (Score: {score:.2f})" for keyword, score in filtered_keywords])

                file_name = "extracted_keywords.txt"

                href = f"data:text/plain;charset=utf-8,{keywords_content.encode('utf-8').decode('latin-1')}"
                st.markdown(f'<a href="{href}" download="{file_name}">Click here to download the keywords file</a>', unsafe_allow_html=True)

            if filtered_keywords:
                col_count = 4
                col_keywords = [[] for _ in range(col_count)]

                for i, (keyword, score) in enumerate(filtered_keywords):
                    col_index = i % col_count
                    col_keywords[col_index].append((keyword, score))

                cols = st.columns(col_count)
                for col_index, col in enumerate(cols):
                    for keyword, score in col_keywords[col_index]:
                        col.write(f"{keyword} (Score: {score:.3f})")
            else:
                st.write("No keywords found.")

if __name__ == "__main__":
    main()